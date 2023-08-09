import numpy as np
import pathlib
cwd = pathlib.Path.cwd()
import pandas
import itertools
import os
import matplotlib.pyplot as plt
import plotly.graph_objects as go
from PIL import Image
import NC
import pickle
from pptx import Presentation
from pptx.util import Inches, Pt
import GeneClass as Gene
import TimeEmbedding as KTA



def main():
    '''
    '''
    maximum_brightness = True
    kmer = 9
    invert_gray = False

    CGRPerChrome(kmer = kmer, maximum_brightness = maximum_brightness, invert_gray = invert_gray)


def GeneTrajectoryPlots():
    '''
    Generates a 2D k-mer Time series plot for the exons and introns of each.

    Probably would be better to use the CGR Generator and have a function pass, but I'm going to be lazy.
    '''
    # left_t = top_t = width_t = height_t = Inches(1)
    # left_p = top_p = Inches(2)
    # height_p = Inches(5)

    # ppt = Presentation()
    # power_point_file = str(cwd / "ChromPerGene.pptx")

    # try:
    #     ppt.save(power_point_file)
    # except PermissionError as e:
    #     print("Please close the power point and retry")
    #     exit()
    # except Exception as e:
    #     print("New Error:")
    #     print(f"{type(e)}")
    #     print(f"{e}")
    #     exit()

    # blank_slie_layout = ppt.slide_layouts[6]
    # comparison_slide_layout = ppt.slide_layouts[4]
    # title_slide_layout = ppt.slide_layouts[0]

    # master_exon_lens = list()
    # master_intron_lens = list()

    # gene_folder = cwd / "GenePerChrome"
    # gene_folder.mkdir(parents = True, exist_ok = True)

    # # pickle_dict: dict = NC.gene_per_chrome()
    # # with open("GenePerChrom.pkl", "wb") as f:
    # #     pickle.dump(pickle_dict, f)

    # with open(str(cwd / "GenePerChrom.pkl"), "rb") as f:
    #     pickle_dict = pickle.load(f)

    # gene: Gene.Gene
    # for gname, gene in pickle_dict.items():
    #     gname_folder = gene_folder / f"{gname}"
    #     gname_folder.mkdir(parents = True, exist_ok = True)
        
    #     exon_folder = gname_folder / "EXON"
    #     intron_folder = gname_folder / "INTRON"
    #     fullseq_folder = gname_folder / "SEQ"

    #     exon_folder.mkdir(parents=True, exist_ok=True)
    #     intron_folder.mkdir(parents=True, exist_ok=True)
    #     fullseq_folder.mkdir(parents=True, exist_ok=True)

    #     with open(str(gname_folder / "Report.txt"), "w+") as txtf:

    #         intron_lens = []
    #         exon_lens = []

    #         # self.ename, self.gname, self.ncibname = ename, gname, ncibname
    #         # self.chrm, self.strand = chrm, strand
    #         gene_data_line = f"{gene.name}\t{gene.ename}\t{gene.gname}\t{gene.ncibname}\nChrome: {gene.chrm}\tStrand: {gene.strand}\n\n"
            
    #         slide = ppt.slides.add_slide(title_slide_layout)
    #         title = slide.shapes.title
    #         title.text = gene_data_line

    #         txtf.write(gene_data_line)

    #         exon_count, intron_count = 0, 0

    #         exon_seq = ""

    #         e: str
    #         i: str

    #         slide = ppt.slides.add_slide(title_slide_layout)
    #         title = slide.shapes.title
    #         title.text = "Exons"

    #         for e in gene.exon_seq:
    #             slide = ppt.slides.add_slide(blank_slie_layout)

    #             exon_count += 1
    #             exon_seq = f"{exon_seq}{e}"

    #             exon_len = len(e)
    #             exon_lens.append(exon_len)
    #             master_exon_lens.append(exon_len)

    #             exon_data_line = f"Exon {exon_count}: {exon_len}\n"
    #             txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
    #             tf = txBox.text_frame
    #             tf.text = exon_data_line

    #             txtf.write(exon_data_line)

    #             if exon_len >= min_exon:
    #                     mer = nucleotide_counter(e.upper(), kmer)
    #                     cgr = chaos_game_representation(mer, kmer)

    #                     if maximum_brightness:
    #                         super_threshold_indices = cgr > 0
    #                         cgr[super_threshold_indices] = 1000


    #                     filepath = exon_folder / f"{gname}_E_{exon_count}_{exon_len}.png"
    #                     if invert_gray:
    #                         plt.imsave(filepath, cgr, cmap = "gray_r")
    #                     else:
    #                         plt.imsave(filepath, cgr, cmap = "gray")
    #                     pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)
    #                     cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
    #                     with open(cgr_filepath, "wb") as c:
    #                         np.save(c, cgr)
                        

    #             # exit()


    #         slide = ppt.slides.add_slide(blank_slie_layout)

    #         exon_full_data_line = f"\nExon Full Sequence: {len(exon_seq)}\n"
    #         exon_data_meta_line = f"Ave Exon Lenght: {int(len(exon_seq) / exon_count)}\tMin Intron Length: {min(exon_lens)}\tMax Intron Length: {max(exon_lens)}\n\n"

    #         txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = "Full Exon Sequence"
            
    #         p = tf.add_paragraph()
    #         p.text = exon_full_data_line

    #         p = tf.add_paragraph()
    #         p.text = exon_data_meta_line

    #         txtf.write(exon_full_data_line)
    #         txtf.write(exon_data_meta_line)

    #         mer = nucleotide_counter(exon_seq, kmer)
    #         cgr = chaos_game_representation(mer, kmer)

    #         if maximum_brightness:
    #             super_threshold_indices = cgr > 0
    #             cgr[super_threshold_indices] = 100

    #         full_exon_seq = filepath = fullseq_folder / f"{gname}_E_Seq_{len(exon_seq)}.png"
    #         if invert_gray:
    #             plt.imsave(filepath, cgr, cmap = "gray_r")
    #         else:
    #             plt.imsave(filepath, cgr, cmap = "gray")
    #         cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
    #         with open(cgr_filepath, "wb") as c:
    #             np.save(c, cgr)
            

    #         pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)

    #         slide = ppt.slides.add_slide(title_slide_layout)
    #         title = slide.shapes.title
    #         title.text = "Introns"

            
    #         intron_seq = ""

    #         for i in gene.intron_seq:

    #             slide = ppt.slides.add_slide(blank_slie_layout)

    #             intron_count += 1
    #             intron_seq = f"{intron_seq}{i}"


    #             intron_len = len(i)
    #             intron_lens.append(intron_len)
    #             master_intron_lens.append(intron_len)

    #             intron_data_line = f"Intron {intron_count}: {intron_len}\n"

    #             exon_data_line = f"Intron {intron_count}: {intron_len}\n"
    #             txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
    #             tf = txBox.text_frame
    #             tf.text = exon_data_line

    #             txtf.write(intron_data_line)
                
    #             mer = nucleotide_counter(i.upper(), kmer)
    #             cgr = chaos_game_representation(mer, kmer)

    #             if maximum_brightness:
    #                 super_threshold_indices = cgr > 0
    #                 cgr[super_threshold_indices] = 100

    #             filepath = intron_folder / f"{gname}_I_{intron_count}_{intron_len}.png"
    #             if invert_gray:
    #                 plt.imsave(filepath, cgr, cmap = "gray_r")
    #             else:
    #                 plt.imsave(filepath, cgr, cmap = "gray")

    #             cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
    #             with open(cgr_filepath, "wb") as c:
    #                 np.save(c, cgr)
                

    #             pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)


    #         intron_data_meta_line = f"\nAve Intron Lenght: {int(sum(intron_lens) / intron_count)}\tMin Exon Length: {min(intron_lens)}\tMax Exon Length: {max(intron_lens)}\n\n"

    #         slide = ppt.slides.add_slide(blank_slie_layout)
    #         txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = "Combined Intron Sequence"
    #         # tf.text = intron_data_meta_line

    #         p = tf.add_paragraph()
    #         txtf.write(intron_data_meta_line)

    #         mer = nucleotide_counter(intron_seq, kmer)
    #         cgr = chaos_game_representation(mer, kmer)

    #         if maximum_brightness:
    #             super_threshold_indices = cgr > 0
    #             cgr[super_threshold_indices] = 100

    #         full_intron_seq = filepath = fullseq_folder / f"{gname}_I_Seq_{len(intron_seq)}.png"
    #         if invert_gray:
    #             plt.imsave(filepath, cgr, cmap = "gray_r")
    #         else:
    #             plt.imsave(filepath, cgr, cmap = "gray")

    #         cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
    #         with open(cgr_filepath, "wb") as c:
    #             np.save(c, cgr)
            

    #         pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)

    #         gene_full_len_line = f"Full Sequence: {len(gene.full_seq[0].upper())}"

    #         slide = ppt.slides.add_slide(blank_slie_layout)
    #         txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = gene_full_len_line

    #         txtf.write(gene_full_len_line)
    
    #         mer = nucleotide_counter(gene.full_seq[0].upper(), kmer)
    #         cgr = chaos_game_representation(mer, kmer)
    #         if maximum_brightness:
    #             super_threshold_indices = cgr > 0
    #             cgr[super_threshold_indices] = 100

    #         full_gene_seq = filepath = fullseq_folder / f"{gname}_F_Seq_{len(gene.full_seq[0].upper())}.png"
    #         if invert_gray:
    #             plt.imsave(filepath, cgr, cmap = "gray_r")
    #         else:
    #             plt.imsave(filepath, cgr, cmap = "gray")

    #         cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
    #         with open(cgr_filepath, "wb") as c:
    #             np.save(c, cgr)
            
            
    #         pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)

    #         slide = ppt.slides.add_slide(blank_slie_layout)
    #         txBox = slide.shapes.add_textbox(left_t, Inches(0.5), width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = "Full Gene Sequence"
    #         pic = slide.shapes.add_picture(str(full_gene_seq), Inches(0.5), top_p, height = Inches(3.5))

    #         txBox = slide.shapes.add_textbox(Inches(5), Inches(0.5), width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = "Full Exon Sequence"
    #         pic = slide.shapes.add_picture(str(full_exon_seq), Inches(5), top_p, height = Inches(3.5))


    #         slide = ppt.slides.add_slide(blank_slie_layout)
    #         txBox = slide.shapes.add_textbox(left_t, Inches(0.5), width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = "Full Intron Sequence"
    #         pic = slide.shapes.add_picture(str(full_intron_seq), Inches(0.5), top_p, height = Inches(3.5))

    #         txBox = slide.shapes.add_textbox(Inches(5), Inches(0.5), width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = "Full Exon Sequence"
    #         pic = slide.shapes.add_picture(str(full_exon_seq), Inches(5), top_p, height = Inches(3.5))


    #         slide = ppt.slides.add_slide(blank_slie_layout)
    #         txBox = slide.shapes.add_textbox(left_t, Inches(0.5), width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = "Full Gene Sequence"
    #         pic = slide.shapes.add_picture(str(full_gene_seq), Inches(0.5), top_p, height = Inches(3.5))

    #         txBox = slide.shapes.add_textbox(Inches(5), Inches(0.5), width_t, height_t)
    #         tf = txBox.text_frame
    #         tf.text = "Full Intron Sequence"
    #         pic = slide.shapes.add_picture(str(full_intron_seq), Inches(5), top_p, height = Inches(3.5))

    #     # exit()

    # fig = go.Figure()
    # fig.add_trace(go.Histogram(x = master_exon_lens, name = "exon"))
    # fig.add_trace(go.Histogram(x = master_intron_lens, name = "intron"))

    # # fig.show()
    # fig.write_image(str(cwd / "ChromPerGene.png"))
    # slide = ppt.slides.add_slide(blank_slie_layout)
    # pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)

    # ppt.save(str(cwd / "ChromPerGene.pptx"))




def CGRPerChrome(kmer = 6, maximum_brightness = False, invert_gray = False, min_exon = 10, image_type = "CGR", *args, **kwargs):
    '''
    Loads up and outputs a gene from every chromosome and generates CGR for each of its exons and introns.

    There also get put into a PPTX presentation, because fuck doing that by hand.

    Ways that images can be generated:
        CGR = chaos game representation
        KTA = K-mer time analysis
    '''

    left_t = top_t = width_t = height_t = Inches(1)
    left_p = top_p = Inches(2)
    height_p = Inches(5)

    ppt = Presentation()
    power_point_file = str(cwd / "ChromPerGene.pptx")

    try:
        ppt.save(power_point_file)
    except PermissionError as e:
        print("Please close the power point and retry")
        exit()
    except Exception as e:
        print("New Error:")
        print(f"{type(e)}")
        print(f"{e}")
        exit()

    blank_slie_layout = ppt.slide_layouts[6]
    comparison_slide_layout = ppt.slide_layouts[4]
    title_slide_layout = ppt.slide_layouts[0]

    master_exon_lens = list()
    master_intron_lens = list()

    gene_folder = cwd / "GenePerChrome"
    gene_folder.mkdir(parents = True, exist_ok = True)

    # pickle_dict: dict = NC.gene_per_chrome()
    # with open("GenePerChrom.pkl", "wb") as f:
    #     pickle.dump(pickle_dict, f)

    with open(str(cwd / "GenePerChrom.pkl"), "rb") as f:
        pickle_dict = pickle.load(f)

    gene: Gene.Gene
    for gname, gene in pickle_dict.items():
        gname_folder = gene_folder / f"{gname}"
        gname_folder.mkdir(parents = True, exist_ok = True)
        
        exon_folder = gname_folder / "EXON"
        intron_folder = gname_folder / "INTRON"
        fullseq_folder = gname_folder / "SEQ"

        exon_folder.mkdir(parents=True, exist_ok=True)
        intron_folder.mkdir(parents=True, exist_ok=True)
        fullseq_folder.mkdir(parents=True, exist_ok=True)

        with open(str(gname_folder / "Report.txt"), "w+") as txtf:

            intron_lens = []
            exon_lens = []

            # self.ename, self.gname, self.ncibname = ename, gname, ncibname
            # self.chrm, self.strand = chrm, strand
            gene_data_line = f"{gene.name}\t{gene.ename}\t{gene.gname}\t{gene.ncibname}\nChrome: {gene.chrm}\tStrand: {gene.strand}\n\n"
            
            slide = ppt.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = gene_data_line

            txtf.write(gene_data_line)

            exon_count, intron_count = 0, 0

            exon_seq = ""

            e: str
            i: str

            slide = ppt.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = "Exons"

            for e in gene.exon_seq:
                slide = ppt.slides.add_slide(blank_slie_layout)

                exon_count += 1
                exon_seq = f"{exon_seq}{e}"

                exon_len = len(e)
                exon_lens.append(exon_len)
                master_exon_lens.append(exon_len)

                exon_data_line = f"Exon {exon_count}: {exon_len}\n"
                txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
                tf = txBox.text_frame
                tf.text = exon_data_line

                txtf.write(exon_data_line)

                if exon_len >= min_exon:
                        if image_type in "CGR":
                            mer = nucleotide_counter(e.upper(), kmer)
                            plot_data = chaos_game_representation(mer, kmer)
                        elif image_type in "KTA":
                            xy = KTA.time_embedding_v3(e, *args, **kwargs)
                            plot_data = 0

                        if maximum_brightness:
                            super_threshold_indices = plot_data > 0
                            plot_data[super_threshold_indices] = 1000


                        filepath = exon_folder / f"{gname}_E_{exon_count}_{exon_len}.png"
                        if invert_gray:
                            plt.imsave(filepath, cgr, cmap = "gray_r")
                        else:
                            plt.imsave(filepath, cgr, cmap = "gray")
                        pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)
                        cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
                        with open(cgr_filepath, "wb") as c:
                            np.save(c, cgr)
                        

                # exit()


            slide = ppt.slides.add_slide(blank_slie_layout)

            exon_full_data_line = f"\nExon Full Sequence: {len(exon_seq)}\n"
            exon_data_meta_line = f"Ave Exon Lenght: {int(len(exon_seq) / exon_count)}\tMin Intron Length: {min(exon_lens)}\tMax Intron Length: {max(exon_lens)}\n\n"

            txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
            tf = txBox.text_frame
            tf.text = "Full Exon Sequence"
            
            p = tf.add_paragraph()
            p.text = exon_full_data_line

            p = tf.add_paragraph()
            p.text = exon_data_meta_line

            txtf.write(exon_full_data_line)
            txtf.write(exon_data_meta_line)

            mer = nucleotide_counter(exon_seq, kmer)
            cgr = chaos_game_representation(mer, kmer)

            if maximum_brightness:
                super_threshold_indices = cgr > 0
                cgr[super_threshold_indices] = 100

            full_exon_seq = filepath = fullseq_folder / f"{gname}_E_Seq_{len(exon_seq)}.png"
            if invert_gray:
                plt.imsave(filepath, cgr, cmap = "gray_r")
            else:
                plt.imsave(filepath, cgr, cmap = "gray")
            cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
            with open(cgr_filepath, "wb") as c:
                np.save(c, cgr)
            

            pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)

            slide = ppt.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = "Introns"

            
            intron_seq = ""

            for i in gene.intron_seq:

                slide = ppt.slides.add_slide(blank_slie_layout)

                intron_count += 1
                intron_seq = f"{intron_seq}{i}"


                intron_len = len(i)
                intron_lens.append(intron_len)
                master_intron_lens.append(intron_len)

                intron_data_line = f"Intron {intron_count}: {intron_len}\n"

                exon_data_line = f"Intron {intron_count}: {intron_len}\n"
                txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
                tf = txBox.text_frame
                tf.text = exon_data_line

                txtf.write(intron_data_line)
                
                mer = nucleotide_counter(i.upper(), kmer)
                cgr = chaos_game_representation(mer, kmer)

                if maximum_brightness:
                    super_threshold_indices = cgr > 0
                    cgr[super_threshold_indices] = 100

                filepath = intron_folder / f"{gname}_I_{intron_count}_{intron_len}.png"
                if invert_gray:
                    plt.imsave(filepath, cgr, cmap = "gray_r")
                else:
                    plt.imsave(filepath, cgr, cmap = "gray")

                cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
                with open(cgr_filepath, "wb") as c:
                    np.save(c, cgr)
                

                pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)


            intron_data_meta_line = f"\nAve Intron Lenght: {int(sum(intron_lens) / intron_count)}\tMin Exon Length: {min(intron_lens)}\tMax Exon Length: {max(intron_lens)}\n\n"

            slide = ppt.slides.add_slide(blank_slie_layout)
            txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
            tf = txBox.text_frame
            tf.text = "Combined Intron Sequence"
            # tf.text = intron_data_meta_line

            p = tf.add_paragraph()
            txtf.write(intron_data_meta_line)

            mer = nucleotide_counter(intron_seq, kmer)
            cgr = chaos_game_representation(mer, kmer)

            if maximum_brightness:
                super_threshold_indices = cgr > 0
                cgr[super_threshold_indices] = 100

            full_intron_seq = filepath = fullseq_folder / f"{gname}_I_Seq_{len(intron_seq)}.png"
            if invert_gray:
                plt.imsave(filepath, cgr, cmap = "gray_r")
            else:
                plt.imsave(filepath, cgr, cmap = "gray")

            cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
            with open(cgr_filepath, "wb") as c:
                np.save(c, cgr)
            

            pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)

            gene_full_len_line = f"Full Sequence: {len(gene.full_seq[0].upper())}"

            slide = ppt.slides.add_slide(blank_slie_layout)
            txBox = slide.shapes.add_textbox(left_t, top_t, width_t, height_t)
            tf = txBox.text_frame
            tf.text = gene_full_len_line

            txtf.write(gene_full_len_line)
    
            mer = nucleotide_counter(gene.full_seq[0].upper(), kmer)
            cgr = chaos_game_representation(mer, kmer)
            if maximum_brightness:
                super_threshold_indices = cgr > 0
                cgr[super_threshold_indices] = 100

            full_gene_seq = filepath = fullseq_folder / f"{gname}_F_Seq_{len(gene.full_seq[0].upper())}.png"
            if invert_gray:
                plt.imsave(filepath, cgr, cmap = "gray_r")
            else:
                plt.imsave(filepath, cgr, cmap = "gray")

            cgr_filepath = pathlib.Path(os.path.splitext(filepath)[0]).with_suffix(".npy")
            with open(cgr_filepath, "wb") as c:
                np.save(c, cgr)
            
            
            pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)

            slide = ppt.slides.add_slide(blank_slie_layout)
            txBox = slide.shapes.add_textbox(left_t, Inches(0.5), width_t, height_t)
            tf = txBox.text_frame
            tf.text = "Full Gene Sequence"
            pic = slide.shapes.add_picture(str(full_gene_seq), Inches(0.5), top_p, height = Inches(3.5))

            txBox = slide.shapes.add_textbox(Inches(5), Inches(0.5), width_t, height_t)
            tf = txBox.text_frame
            tf.text = "Full Exon Sequence"
            pic = slide.shapes.add_picture(str(full_exon_seq), Inches(5), top_p, height = Inches(3.5))


            slide = ppt.slides.add_slide(blank_slie_layout)
            txBox = slide.shapes.add_textbox(left_t, Inches(0.5), width_t, height_t)
            tf = txBox.text_frame
            tf.text = "Full Intron Sequence"
            pic = slide.shapes.add_picture(str(full_intron_seq), Inches(0.5), top_p, height = Inches(3.5))

            txBox = slide.shapes.add_textbox(Inches(5), Inches(0.5), width_t, height_t)
            tf = txBox.text_frame
            tf.text = "Full Exon Sequence"
            pic = slide.shapes.add_picture(str(full_exon_seq), Inches(5), top_p, height = Inches(3.5))


            slide = ppt.slides.add_slide(blank_slie_layout)
            txBox = slide.shapes.add_textbox(left_t, Inches(0.5), width_t, height_t)
            tf = txBox.text_frame
            tf.text = "Full Gene Sequence"
            pic = slide.shapes.add_picture(str(full_gene_seq), Inches(0.5), top_p, height = Inches(3.5))

            txBox = slide.shapes.add_textbox(Inches(5), Inches(0.5), width_t, height_t)
            tf = txBox.text_frame
            tf.text = "Full Intron Sequence"
            pic = slide.shapes.add_picture(str(full_intron_seq), Inches(5), top_p, height = Inches(3.5))

        # exit()

    fig = go.Figure()
    fig.add_trace(go.Histogram(x = master_exon_lens, name = "exon"))
    fig.add_trace(go.Histogram(x = master_intron_lens, name = "intron"))

    # fig.show()
    fig.write_image(str(cwd / "ChromPerGene.png"))
    slide = ppt.slides.add_slide(blank_slie_layout)
    pic = slide.shapes.add_picture(str(filepath), left_p, top_p, height = height_p)

    ppt.save(str(cwd / "ChromPerGene.pptx"))


# exon_images = cwd / "FractalImageEvI_SL" / "EXON"
# intron_images = cwd / "FractalImageEvI_SL" / "INTRON"
# # utr5_images = cwd / "FractalImage" / "UTR5"
# # utr3_images = cwd / "FractalImage" / "UTR3"
# # exon_images = cwd / "FractalImage" / "Exon"

# exon_images.mkdir(parents = True, exist_ok = True)
# intron_images.mkdir(parents = True, exist_ok = True)
# # utr5_images.mkdir(parents = True, exist_ok = True)
# # utr3_images.mkdir(parents = True, exist_ok = True)
# # exon_images.mkdir(parents = True, exist_ok = True)

# kmer = 6

def manhattan_position(nuc: int, x0: int):
    '''
    '''
    x1: np.array = (nuc - x0) / 2

    return x1.astype(int)


def nucleotide_permutations(sequence: str = "ACGT", length: int = 3) -> dict:
    nuc_perm = dict()

    if len(sequence) < length:
        return None

    perms = itertools.permutations(sequence, length)
    for p in perms:

        key = ""
        for n in p:
            key = f"{key}{n}"

        nuc_perm[key] = 0

    return nuc_perm


def nucleotide_counter(sequence: str, window_size: int):
    '''
    '''
    keys: set = set()
    counter = dict()
    master_count = 0


    for i in range(len(sequence) - window_size):
        seq = sequence[i: i + window_size]

        if seq not in keys:
            keys.add(seq)
            counter[seq] = 1
            master_count += 1

        else:
            counter[seq] += 1
            master_count += 1

    # for key, value in counter.items():
    #     counter[key] = value / master_count

    return counter



def correlation_dimension_figure(cgr):
    '''
    '''

    local_c: dict = correlation_dimension(cgr)
    
    fig = go.Figure()
    fig.add_trace(go.Scatter(x = list(local_c.keys()), y = list(local_c.values())))
    fig.show()




def correlation_dimension(cgr: np.ndarray) -> dict:
    '''
    Take the euclidean distance between every point in the matrix. Start with a small radius then increase this to the whole thing.

    Will need to find the size of the input matrix. Assume it's always a square (in this case it is).
    Also don't need to calculate every single point over and over again. Can find a way to calculate only half the values. Also once the values
    are found, save them so I don't need to recalculate them.

    Also using the Heaviside function, so if (xi - xj) is a point: 1. Else 0.
    '''
    non_zero_indices = np.transpose(np.where(cgr > 0))
    hash_index = lambda x: f"{x[0]}, {x[1]}"

    C = dict()

    distances = dict()
    for i in non_zero_indices:
        hash_i = hash_index(i)
        distances[hash_i] = dict()
        for j in non_zero_indices:
            hash_j = hash_index(j)
            if hash_i not in hash_j:
                distances[hash_i][hash_j] = 0

    l, w = cgr.shape
    # N = (l*w)**(-2)
    N = non_zero_indices.shape[0]**(-2)
    size = int(np.sqrt(l**2 + w**2)) + 1

    for e in range(size):
        C[e] = 0
        
        for i in non_zero_indices:
            hash_i = hash_index(i)

            for j in non_zero_indices:
                hash_j = hash_index(j)

                if hash_i not in hash_j:
                
                    if distances[hash_i][hash_j] == 0:
                        d = heaviside_function(i, j, e)
                        distances[hash_i][hash_j] = d

                    C[e] += distances[hash_i][hash_j]


        C[e] = C[e] * N  # note the multiplication: all the work for making N be 1/N**2 was taken care of earlier

    # exit()

    return C

    # for r in range(size):
    #     for index in non_zero_indices:
    #         c = heaviside_function()



def heaviside_function(xyi: float, xyj: float, r: float) -> int:
    '''
    Actually computes the manhattan distance but that's because I'm lazy and don't want to deal with sqrt right now.
    '''

    # d = abs(xyi[0] - xyj[0]) + abs(xyi[1] - xyj[1])
    d = np.sqrt(np.power(xyi[0] - xyj[0], 2) + np.power(xyi[1] - xyj[1], 2))

    if ((r - d) >= 0):
        return 1
    else:
        return 0


def chaos_game_representation(probabilities: dict, k) -> np.ndarray:
    '''
    '''
    array_size = int(np.sqrt(4**k))

    cgr = np.zeros(shape = (array_size, array_size))

    for key, value in probabilities.items():
        maxx = array_size
        maxy = array_size
        posx = 1
        posy = 1
 
        for char in key:
            if char == "T":
                posx += maxx / 2
            elif char == "C":
                posy += maxy / 2
            elif char == "G":
                posx += maxx / 2
                posy += maxy / 2
            maxx /=  2
            maxy /= 2

        # print(int(posy), int(posx))
        cgr[int(posy - 1)][int(posx - 1)] = value


    return cgr


def gif_generator(path: str, filename: str):
    '''
    '''
    png_files = list(path.rglob("*.png"))
    png_files.sort(key = lambda x: os.path.getmtime(x))

    frames = []
    for png in png_files:
        new_frame = Image.open(png)
        # print(png)
        frames.append(new_frame)


    frames[0].save(filename, format="GIF", append_images=frames[1:], save_all = True, loop = 0, duration = 225) 


# train_data: pandas.DataFrame = pandas.read_pickle(cwd / "TrainingData_SameSize.pkl")
# # print(train_data.shape)

# keep = np.where(train_data["Seq"].str.len() >= 100)[0]
# # print(keep)
# train_data = train_data.iloc[keep, :]

# rows, cols = train_data.shape

# train_data = train_data.reset_index()

# # print(train_data["Type"].unique())

# # print(train_data.shape)
# # print(train_data)

# print(train_data)



# exon = 0
# intron, intname = 0, "Intron"
# cds, exname = 0, "Exon"
# # utr3, ut3name = 0, "UTR3"
# # utr5, ut5name = 0, "UTR5"

# for row in range(rows):
#     seq = train_data.loc[row, "Seq"]
#     typ = train_data.loc[row, "Type"]

#     mer = nucleotide_counter(seq, kmer)
#     cgr = chaos_game_representation(mer, kmer)

#     # print(cgr)
#     # print(np.max(cgr), np.min(cgr))

#     plt.imshow(cgr, cmap = "gray")

#     if typ in intname:
#         filepath = intron_images / f"Intron_{intron}.png"
#         plt.imsave(filepath, cgr, cmap = "gray")
#         intron += 1

#     elif typ in exname:

#         filepath = exon_images / f"Exon_{exon}.png"
#         plt.imsave(filepath, cgr, cmap = "gray")

#         exon += 1

if __name__ in "__main__":
    main()
